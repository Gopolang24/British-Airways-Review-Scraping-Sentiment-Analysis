from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import pandas as pd

# Load dataset
csv_path = "data/reviews.csv"
df = pd.read_csv(csv_path)

# Ensure column names are properly formatted
df.columns = df.columns.str.strip().str.lower()

# Create a new presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Airline Review Scraping Summary"
subtitle = slide.placeholders[1]
subtitle.text = "Key Findings from Data Analysis"

# Slide 2: Dataset Overview
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Dataset Overview"
content = slide.placeholders[1]
content.text = (
    "- Reviews include ratings for seat comfort, cabin staff, food, and entertainment.\n"
    "- Data categorized by seat type, traveler type, and route.\n"
    "- Helps analyze customer satisfaction trends."
)

# Check if 'sentiment_label' column exists
if 'sentiment_label' in df.columns:
    # Generate Sentiment Distribution Visualization
    sentiment_counts = df['sentiment_label'].value_counts()
    plt.figure(figsize=(6,4))
    sentiment_counts.plot(kind='bar', color=['green', 'gray', 'red'])
    plt.title('Sentiment Distribution')
    plt.xlabel('Sentiment')
    plt.ylabel('Number of Reviews')
    plt.xticks(rotation=0)
    sentiment_chart = "data/sentiment_chart.png"
    plt.savefig(sentiment_chart)
    plt.close()

    # Slide 3: Sentiment Analysis with Chart
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Sentiment Analysis"
    content = slide.placeholders[1]
    content.text = (
        "- Reviews classified as Positive, Neutral, or Negative.\n"
        "- Text processing techniques used for classification.\n"
        "- Provides insights into passenger experience."
    )
    slide.shapes.add_picture(sentiment_chart, Inches(1), Inches(2.5), width=Inches(6))
else:
    print("Warning: 'sentiment_label' column not found in dataset.")

# Check if 'seat type' and 'overall rating' columns exist
if 'seat type' in df.columns and 'overall rating' in df.columns:
    # Generate Key Trends Visualization
    class_ratings = df.groupby('seat type')['overall rating'].mean()
    plt.figure(figsize=(6,4))
    class_ratings.plot(kind='bar', color=['blue', 'orange', 'purple'])
    plt.title('Average Ratings by Seat Type')
    plt.xlabel('Seat Type')
    plt.ylabel('Average Rating')
    plt.xticks(rotation=0)
    ratings_chart = "data/ratings_chart.png"
    plt.savefig(ratings_chart)
    plt.close()

    # Slide 4: Key Trends with Chart
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Trends"
    content = slide.placeholders[1]
    content.text = (
        "- Premium Economy & Business Class: Higher satisfaction.\n"
        "- Economy Class: More negative reviews on food & entertainment.\n"
        "- Long international routes show higher satisfaction."
    )
    slide.shapes.add_picture(ratings_chart, Inches(1), Inches(2.5), width=Inches(6))
else:
    print("Warning: 'seat type' or 'overall rating' column not found in dataset.")

# Slide 5: Scraping Methodology
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Scraping Methodology"
content = slide.placeholders[1]
content.text = (
    "- Data extracted using BeautifulSoup.\n"
    "- Ratings, reviews, and traveler details parsed from website tables.\n"
    "- Ensured structured data collection for analysis."
)

# Save the PowerPoint file
pptx_path = "data/Presentation Template-Task 1.pptx"
prs.save(pptx_path)

# Provide download link
pptx_path